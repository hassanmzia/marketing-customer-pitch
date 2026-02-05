#!/usr/bin/env python3
"""Generate a professional PowerPoint presentation for AI Marketing Customer Pitch Assistant."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0B, 0x1D, 0x3A)
MID_BLUE    = RGBColor(0x14, 0x3A, 0x6B)
TEAL        = RGBColor(0x0D, 0x7C, 0x85)
LIGHT_TEAL  = RGBColor(0x17, 0xA2, 0xAD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xE8, 0xEC, 0xF1)
ACCENT_GOLD = RGBColor(0xF0, 0xC0, 0x40)
SOFT_WHITE  = RGBColor(0xF5, 0xF7, 0xFA)
BORDER_TEAL = RGBColor(0x0A, 0x5E, 0x66)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ── Helpers ────────────────────────────────────────────────────────────────
def _add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def _set_text(shape, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def _add_bullet_para(tf, text, font_size=14, color=WHITE, bold=False, level=0, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_before = space_before
    return p


def _title_bar(slide, title_text):
    """Add a consistent title bar at the top of content slides."""
    bar = _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), MID_BLUE)
    _set_text(bar, title_text, font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    bar.text_frame.paragraphs[0].space_before = Pt(8)
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.margin_top = Inches(0.15)
    # accent line
    _add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.05), LIGHT_TEAL)


def _arrow_right(slide, left, top, width=Inches(0.6), height=Inches(0.35), color=LIGHT_TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_add_bg(slide1, DARK_BLUE)

# Decorative top accent
_add_rect(slide1, Inches(0), Inches(0), SLIDE_W, Inches(0.08), LIGHT_TEAL)

# Title
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(1.6))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AI Marketing Customer Pitch Assistant"
p.font.size = Pt(44)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Accent line under title
_add_rect(slide1, Inches(4.5), Inches(3.7), Inches(4.3), Inches(0.06), LIGHT_TEAL)

# Subtitle
sub_box = slide1.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(1.0))
tf = sub_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Technical Architecture Overview"
p.font.size = Pt(28)
p.font.color.rgb = LIGHT_TEAL
p.font.bold = False
p.alignment = PP_ALIGN.CENTER

# Bottom decorative bar
_add_rect(slide1, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), TEAL)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – System Architecture (5 layers)
# ══════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide2, DARK_BLUE)
_title_bar(slide2, "System Architecture")

layers = [
    ("1  CLIENT LAYER",
     "React 18  +  TypeScript  +  Vite  +  Tailwind CSS  +  React Query  +  Zustand",
     RGBColor(0x14, 0x6E, 0x8A)),
    ("2  GATEWAY LAYER",
     "Nginx (reverse proxy, static assets, gzip, security headers)  +  Express BFF (rate limiting, CORS, WebSocket)",
     RGBColor(0x0F, 0x5E, 0x70)),
    ("3  BACKEND LAYER",
     "Django 5.1  +  DRF  (6 apps: core, customers, pitches, campaigns, agents, analytics)",
     RGBColor(0x0D, 0x50, 0x63)),
    ("4  AI / AGENT LAYER",
     "FastMCP Server with 10 tools  |  Multi-Agent Pipeline (Research \u2192 Generator \u2192 Scorer \u2192 Refiner)  |  LangChain + OpenAI gpt-4o-mini",
     RGBColor(0x0A, 0x42, 0x55)),
    ("5  DATA LAYER",
     "PostgreSQL 16  +  Redis 7 (cache, Celery broker, WebSocket channels)",
     RGBColor(0x08, 0x35, 0x48)),
]

card_left = Inches(0.6)
card_w = Inches(12.1)
card_h = Inches(1.0)
start_top = Inches(1.45)
gap = Inches(0.18)

for i, (label, desc, bg_c) in enumerate(layers):
    top = start_top + i * (card_h + gap)
    card = _add_rounded_rect(slide2, card_left, top, card_w, card_h, bg_c, BORDER_TEAL)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = ACCENT_GOLD
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(4)

    # Down arrow between layers (except after last)
    if i < len(layers) - 1:
        arrow_shape = slide2.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(6.4), top + card_h, Inches(0.5), Inches(0.16), 
        )
        arrow_shape.fill.solid()
        arrow_shape.fill.fore_color.rgb = LIGHT_TEAL
        arrow_shape.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Multi-Agent Pipeline
# ══════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide3, DARK_BLUE)
_title_bar(slide3, "AI Multi-Agent Pipeline (A2A Communication)")

steps = [
    ("Orchestrator",          "Receives pitch request\nand coordinates agents", RGBColor(0x0D, 0x7C, 0x85)),
    ("Research Agent",        "Gathers customer\nintelligence & context",       RGBColor(0x10, 0x6B, 0x78)),
    ("Pitch Generator",      "Creates personalized\npitch content",            RGBColor(0x13, 0x5A, 0x6B)),
    ("Scorer Agent",          "Evaluates persuasiveness,\nclarity, relevance (0\u20131)", RGBColor(0x16, 0x49, 0x5E)),
]

box_w = Inches(2.4)
box_h = Inches(1.5)
total_w = 4 * box_w.inches + 3 * 0.7
start_x = (13.333 - total_w) / 2
y_top = Inches(1.6)

for i, (name, desc, bg) in enumerate(steps):
    lx = Inches(start_x + i * (box_w.inches + 0.7))
    box = _add_rounded_rect(slide3, lx, y_top, box_w, box_h, bg, BORDER_TEAL)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

    if i < len(steps) - 1:
        ax = Inches(start_x + (i + 1) * box_w.inches + i * 0.7 + 0.05)
        _arrow_right(slide3, ax, Inches(2.15), Inches(0.6), Inches(0.35))

# Decision diamond area
decision_top = Inches(3.5)
dec_box = _add_rounded_rect(slide3, Inches(3.5), decision_top, Inches(6.3), Inches(1.2),
                            RGBColor(0x0A, 0x42, 0x55), LIGHT_TEAL)
tf = dec_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "SCORING DECISION"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD
p.alignment = PP_ALIGN.CENTER
for line in [
    "Score \u2265 0.7  \u2192  APPROVED  \u2013 pitch delivered to user",
    "Score < 0.7  \u2192  REFINER AGENT  \u2013 iterative improvement (max 3 iterations)",
]:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(13)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

# Output box
out_box = _add_rounded_rect(slide3, Inches(2.5), Inches(5.1), Inches(8.3), Inches(1.1),
                            MID_BLUE, BORDER_TEAL)
tf = out_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.08)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "FINAL OUTPUT"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Final pitch delivered with full audit trail  \u2022  Version history  \u2022  A2A message trace"
p2.font.size = Pt(13)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(4)

# Arrow between decision and output
slide3.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW,
    Inches(6.4), Inches(4.72), Inches(0.5), Inches(0.35),
).fill.solid()
slide3.shapes[-1].fill.fore_color.rgb = LIGHT_TEAL
slide3.shapes[-1].line.fill.background()


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Technology Stack (two columns)
# ══════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide4, DARK_BLUE)
_title_bar(slide4, "Technology Stack")

col_w = Inches(5.8)
col_h = Inches(5.4)
col_y = Inches(1.5)

# Left column – Frontend
left_col = _add_rounded_rect(slide4, Inches(0.6), col_y, col_w, col_h,
                              RGBColor(0x10, 0x2E, 0x50), BORDER_TEAL)
tf = left_col.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "FRONTEND"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD
p.alignment = PP_ALIGN.CENTER

frontend_items = [
    "React 18", "TypeScript", "Vite", "Tailwind CSS",
    "React Query", "Zustand", "Recharts", "Framer Motion"
]
for item in frontend_items:
    _add_bullet_para(tf, f"\u2022  {item}", font_size=16, color=WHITE, space_before=Pt(8))

# Right column – Backend
right_col = _add_rounded_rect(slide4, Inches(6.9), col_y, col_w, col_h,
                               RGBColor(0x10, 0x2E, 0x50), BORDER_TEAL)
tf = right_col.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "BACKEND"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD
p.alignment = PP_ALIGN.CENTER

backend_items = [
    "Django 5.1", "Django REST Framework", "Celery 5.4",
    "PostgreSQL 16", "Redis 7", "LangChain", "FastMCP", "OpenAI API"
]
for item in backend_items:
    _add_bullet_para(tf, f"\u2022  {item}", font_size=16, color=WHITE, space_before=Pt(8))


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Docker Container Architecture
# ══════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide5, DARK_BLUE)
_title_bar(slide5, "Docker Container Architecture (10 Services)")

services = [
    ("postgres",       ":5432", "PostgreSQL 16 database"),
    ("redis",          ":6379", "Cache, Celery broker, WebSocket channels"),
    ("backend",        ":8064", "Django 5.1 + DRF application"),
    ("mcp-server",     ":8165", "FastMCP AI tool server"),
    ("celery-worker",  "3 queues", "Async task processing"),
    ("celery-beat",    "scheduler", "Periodic task scheduling"),
    ("flower",         ":5555", "Celery monitoring dashboard"),
    ("bff",            ":4064", "Express.js Backend-for-Frontend"),
    ("frontend/nginx", ":3064", "React SPA + reverse proxy"),
]

col_count = 3
row_count = 3
card_w = Inches(3.7)
card_h = Inches(1.4)
x_start = Inches(0.65)
y_start = Inches(1.5)
x_gap = Inches(0.35)
y_gap = Inches(0.3)

for idx, (name, port, desc) in enumerate(services):
    col = idx % col_count
    row = idx // col_count
    lx = x_start + col * (card_w + x_gap)
    ly = y_start + row * (card_h + y_gap)
    card = _add_rounded_rect(slide5, lx, ly, card_w, card_h,
                              RGBColor(0x10, 0x2E, 0x50), BORDER_TEAL)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p2 = tf.add_paragraph()
    p2.text = port
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT_TEAL
    p2.font.bold = True
    p2.space_before = Pt(3)
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(12)
    p3.font.color.rgb = WHITE
    p3.space_before = Pt(3)

# Footer note – use the remaining space beneath the grid
note_box = slide5.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5))
tf = note_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "All services on  pitch-network  bridge  \u2022  Health checks configured on every service"
p.font.size = Pt(14)
p.font.color.rgb = LIGHT_TEAL
p.font.bold = True
p.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Request Data Flow
# ══════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide6, DARK_BLUE)
_title_bar(slide6, "Request Data Flow")

flow_nodes = [
    "User\n(Browser)",
    "Nginx\n:3064",
    "BFF Express\n:4064",
    "Django DRF\n:8064",
    "Celery\nWorker",
    "AgentService\n(LangChain)",
    "MCP\nTools",
    "OpenAI\nAPI",
]

node_w = Inches(1.25)
node_h = Inches(0.95)
total_nodes_w = len(flow_nodes) * node_w.inches + (len(flow_nodes) - 1) * 0.25
sx = (13.333 - total_nodes_w) / 2
ny = Inches(2.0)

for i, label in enumerate(flow_nodes):
    nx = Inches(sx + i * (node_w.inches + 0.25))
    box = _add_rounded_rect(slide6, nx, ny, node_w, node_h,
                             RGBColor(0x0D, 0x50, 0x63), LIGHT_TEAL)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if i < len(flow_nodes) - 1:
        ax = Inches(sx + (i + 1) * node_w.inches + i * 0.25 + 0.02)
        _arrow_right(slide6, ax, Inches(2.28), Inches(0.22), Inches(0.3))

# Storage boxes
store_y = Inches(3.8)
pg_box = _add_rounded_rect(slide6, Inches(3.0), store_y, Inches(3.5), Inches(1.2),
                            RGBColor(0x10, 0x2E, 0x50), BORDER_TEAL)
tf = pg_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.1)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "PostgreSQL"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Results stored persistently"
p2.font.size = Pt(12)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(4)

rd_box = _add_rounded_rect(slide6, Inches(7.0), store_y, Inches(3.5), Inches(1.2),
                            RGBColor(0x10, 0x2E, 0x50), BORDER_TEAL)
tf = rd_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.1)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Redis"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Cached results & session data"
p2.font.size = Pt(12)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(4)

# WebSocket note
ws_box = _add_rounded_rect(slide6, Inches(3.5), Inches(5.5), Inches(6.3), Inches(0.9),
                            RGBColor(0x0A, 0x42, 0x55), LIGHT_TEAL)
tf = ws_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Real-time updates delivered via WebSocket (Django Channels + Redis)"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Key Features
# ══════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide7, DARK_BLUE)
_title_bar(slide7, "Key Features")

features = [
    "Customer 360\u00b0 View with AI-enriched data",
    "AI-powered pitch generation with scoring and iterative refinement",
    "Multi-channel campaign management with A/B testing",
    "Agent-to-Agent (A2A) communication with full audit trail",
    "PDF / DOCX / TXT export for pitches",
    "Real-time WebSocket updates",
    "10 specialized MCP tools for AI operations",
    "Comprehensive analytics and agent performance tracking",
]

feat_box = _add_rounded_rect(slide7, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5),
                              RGBColor(0x10, 0x2E, 0x50), BORDER_TEAL)
tf = feat_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.5)
tf.margin_top = Inches(0.35)

for i, feat in enumerate(features):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"\u2714   {feat}"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.space_before = Pt(14)
    # Alternate slight teal highlight on the checkmark via bold
    p.font.bold = False


# ── Save ───────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "technical-architecture.pptx")
prs.save(out_path)
print(f"Presentation saved to {out_path}")
